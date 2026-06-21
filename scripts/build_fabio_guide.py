#!/usr/bin/env python3
"""Genera la guida PowerPoint di CustomerSpeed per Fabio.

Deck 16:9, italiano, brand indigo. Esecuzione:
    python3 scripts/build_fabio_guide.py
Output: Guida-CustomerSpeed-Fabio.pptx (nella root del progetto).
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

# ── Brand ────────────────────────────────────────────────────────────────────
INDIGO       = RGBColor(0x63, 0x66, 0xF1)   # accent
INDIGO_DARK  = RGBColor(0x43, 0x38, 0xCA)
INK          = RGBColor(0x1E, 0x1B, 0x4B)   # titoli
BODY         = RGBColor(0x33, 0x33, 0x3D)   # testo
MUTED        = RGBColor(0x6B, 0x70, 0x80)   # secondario
WHITE        = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT        = RGBColor(0xF4, 0xF4, 0xF8)   # box chiari
FONT         = "Arial"

EMU_W, EMU_H = Inches(13.333), Inches(7.5)

prs = Presentation()
prs.slide_width = EMU_W
prs.slide_height = EMU_H
BLANK = prs.slide_layouts[6]


def _set(run, size, color, bold=False, italic=False, font=FONT):
    run.font.size = Pt(size)
    run.font.color.rgb = color
    run.font.bold = bold
    run.font.italic = italic
    run.font.name = font


def box(slide, x, y, w, h, fill=None, line=None):
    shp = slide.shapes.add_shape(1, x, y, w, h)  # rectangle
    shp.fill.solid()
    if fill is None:
        shp.fill.background()
    else:
        shp.fill.fore_color.rgb = fill
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line
        shp.line.width = Pt(1)
    shp.shadow.inherit = False
    return shp


def text(slide, x, y, w, h, runs, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
         space_after=6, line_spacing=1.05):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    first = True
    for item in runs:
        # item: (text, size, color, bold, italic, level, bullet)
        txt, size, color, bold, italic, level, bullet = item
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.alignment = align
        p.level = level
        p.space_after = Pt(space_after)
        p.line_spacing = line_spacing
        prefix = "•  " if bullet == 1 else ("–  " if bullet == 2 else "")
        r = p.add_run()
        r.text = prefix + txt
        _set(r, size, color, bold, italic)
    return tb


def b(txt, size=16, color=BODY, bold=False, italic=False, level=0, bullet=1):
    return (txt, size, color, bold, italic, level, bullet)


PAGE = {"n": 0}


def chrome(slide, title, kicker=None):
    """Barra superiore + titolo + footer, layout standard delle slide contenuto."""
    box(slide, 0, 0, EMU_W, Inches(1.25), fill=WHITE)
    box(slide, 0, Inches(1.25), EMU_W, Pt(3), fill=INDIGO)
    box(slide, Inches(0.55), Inches(0.42), Inches(0.12), Inches(0.5), fill=INDIGO)
    if kicker:
        text(slide, Inches(0.8), Inches(0.30), Inches(11), Inches(0.3),
             [b(kicker.upper(), 11, INDIGO, bold=True, bullet=0)])
        text(slide, Inches(0.8), Inches(0.56), Inches(11.7), Inches(0.6),
             [b(title, 26, INK, bold=True, bullet=0)])
    else:
        text(slide, Inches(0.8), Inches(0.40), Inches(11.7), Inches(0.7),
             [b(title, 28, INK, bold=True, bullet=0)])
    # footer
    PAGE["n"] += 1
    text(slide, Inches(0.55), Inches(7.02), Inches(8), Inches(0.35),
         [b("CustomerSpeed · Guida rapida per Fabio", 9, MUTED, bullet=0)])
    text(slide, Inches(11.8), Inches(7.02), Inches(1.0), Inches(0.35),
         [b(str(PAGE["n"]), 9, MUTED, bullet=0)], align=PP_ALIGN.RIGHT)


def content_slide(title, kicker, body_runs, note=None):
    s = prs.slides.add_slide(BLANK)
    box(s, 0, 0, EMU_W, EMU_H, fill=WHITE)
    chrome(s, title, kicker)
    text(s, Inches(0.85), Inches(1.6), Inches(11.6), Inches(5.1), body_runs,
         space_after=9, line_spacing=1.08)
    if note:
        box(s, Inches(0.85), Inches(6.05), Inches(11.6), Inches(0.78), fill=LIGHT)
        box(s, Inches(0.85), Inches(6.05), Pt(4), Inches(0.78), fill=INDIGO)
        text(s, Inches(1.05), Inches(6.14), Inches(11.2), Inches(0.62),
             [b(note, 12.5, INK, italic=True, bullet=0)], anchor=MSO_ANCHOR.MIDDLE)
    return s


def two_col_slide(title, kicker, left_head, left_runs, right_head, right_runs):
    s = prs.slides.add_slide(BLANK)
    box(s, 0, 0, EMU_W, EMU_H, fill=WHITE)
    chrome(s, title, kicker)
    # left card
    box(s, Inches(0.85), Inches(1.7), Inches(5.65), Inches(4.9), fill=LIGHT)
    text(s, Inches(1.1), Inches(1.95), Inches(5.2), Inches(0.5),
         [b(left_head, 16, INDIGO_DARK, bold=True, bullet=0)])
    text(s, Inches(1.1), Inches(2.55), Inches(5.2), Inches(3.9), left_runs, space_after=8)
    # right card
    box(s, Inches(6.8), Inches(1.7), Inches(5.65), Inches(4.9), fill=LIGHT)
    text(s, Inches(7.05), Inches(1.95), Inches(5.2), Inches(0.5),
         [b(right_head, 16, INDIGO_DARK, bold=True, bullet=0)])
    text(s, Inches(7.05), Inches(2.55), Inches(5.2), Inches(3.9), right_runs, space_after=8)
    return s


# ── 1. Cover ─────────────────────────────────────────────────────────────────
s = prs.slides.add_slide(BLANK)
box(s, 0, 0, EMU_W, EMU_H, fill=INDIGO)
box(s, 0, 0, EMU_W, EMU_H, fill=INDIGO)
box(s, Inches(0.9), Inches(2.5), Inches(1.4), Pt(6), fill=WHITE)
text(s, Inches(0.85), Inches(1.5), Inches(11), Inches(0.6),
     [b("CUSTOMERSPEED", 18, WHITE, bold=True, bullet=0)])
text(s, Inches(0.85), Inches(2.75), Inches(11.6), Inches(2.0),
     [b("Guida rapida all'uso del CRM", 44, WHITE, bold=True, bullet=0)])
text(s, Inches(0.88), Inches(4.35), Inches(11), Inches(0.8),
     [b("Il tuo strumento per gestire lead, pipeline e clienti — passo per passo.",
        18, RGBColor(0xE7, 0xE7, 0xFB), bullet=0)])
text(s, Inches(0.88), Inches(6.3), Inches(11), Inches(0.6),
     [b("Preparata per Fabio · accesso riservato", 13,
        RGBColor(0xCD, 0xCE, 0xF7), bullet=0)])

# ── 2. Indice ────────────────────────────────────────────────────────────────
content_slide(
    "Cosa trovi in questa guida", "Indice",
    [
        b("Accesso e primo login", 17, INK, bold=True),
        b("La schermata principale e la navigazione", 17, INK, bold=True),
        b("Dashboard: i numeri chiave a colpo d'occhio", 17, INK, bold=True),
        b("Lead: elenco, creazione e scheda di dettaglio", 17, INK, bold=True),
        b("Pipeline: la board kanban e gli stati", 17, INK, bold=True),
        b("Appuntamenti, esportazione dati e impostazioni", 17, INK, bold=True),
        b("Privacy, cookie e sicurezza (reCAPTCHA)", 17, INK, bold=True),
    ],
    note="Suggerimento: l'interfaccia è in italiano e funziona su computer, tablet e telefono.")

# ── 3. Accesso ───────────────────────────────────────────────────────────────
content_slide(
    "Accedere a CustomerSpeed", "Primo passo",
    [
        b("Apri il browser e vai all'indirizzo del tuo spazio riservato:", 16, BODY, bullet=0),
        b("https://customer-speed.vercel.app/login?org=fabio", 17, INDIGO_DARK, bold=True, level=1, bullet=2),
        b("Inserisci Email e Password che ti sono state fornite, poi premi “Accedi”.", 16, BODY),
        b("Hai dimenticato la password? Usa “Password dimenticata?” per riceverne una nuova via email.", 16, BODY),
        b("Per uscire in sicurezza usa “Esci” dal menu in alto quando hai finito.", 16, BODY),
    ],
    note="Importante: usa sempre il link con ?org=fabio — è ciò che ti porta nel TUO spazio.")

# 3b. credenziali (card)
s = prs.slides.add_slide(BLANK)
box(s, 0, 0, EMU_W, EMU_H, fill=WHITE)
chrome(s, "Le tue credenziali", "Accesso")
box(s, Inches(0.85), Inches(1.8), Inches(11.6), Inches(2.7), fill=LIGHT)
box(s, Inches(0.85), Inches(1.8), Pt(5), Inches(2.7), fill=INDIGO)
text(s, Inches(1.2), Inches(2.05), Inches(10.8), Inches(2.2),
     [
        b("Indirizzo:  https://customer-speed.vercel.app/login?org=fabio", 17, INK, bold=True, bullet=0),
        b("Email:  _______________________________", 17, BODY, bullet=0),
        b("Password:  ____________________________", 17, BODY, bullet=0),
     ], space_after=16, line_spacing=1.2)
text(s, Inches(0.85), Inches(4.8), Inches(11.6), Inches(1.6),
     [
        b("Conserva queste credenziali in un luogo sicuro e non condividerle.", 15, BODY),
        b("Alla prima occasione, cambia la password con una personale e robusta.", 15, BODY),
     ], space_after=10)

# ── 4. Navigazione ───────────────────────────────────────────────────────────
two_col_slide(
    "La schermata principale", "Orientarsi",
    "Il menu laterale", [
        b("Dashboard — riepilogo e numeri chiave", 15, BODY),
        b("Pipeline — la board dei tuoi lead per stato", 15, BODY),
        b("Lead — l'elenco completo dei contatti", 15, BODY),
        b("Appuntamenti — la tua agenda", 15, BODY),
        b("Impostazioni — preferenze e aspetto", 15, BODY),
    ],
    "Comandi utili", [
        b("Tema chiaro / scuro — interruttore in alto", 15, BODY),
        b("Su telefono il menu si chiude e si apre con l'icona ☰", 15, BODY),
        b("Il nome in alto a destra apre il menu profilo e “Esci”", 15, BODY),
        b("Tutto è navigabile anche da tastiera", 15, BODY),
    ])

# ── 5. Dashboard ─────────────────────────────────────────────────────────────
content_slide(
    "Dashboard: i numeri a colpo d'occhio", "Panoramica",
    [
        b("È la prima pagina dopo il login. Mostra lo stato del tuo lavoro in sintesi:", 16, BODY, bullet=0),
        b("Quanti lead hai in totale e quanti attivi", 16, BODY),
        b("Lead vinti e persi (conversioni)", 16, BODY),
        b("Distribuzione dei lead nei vari stati della pipeline", 16, BODY),
        b("Valore economico generato (fatturato dai lead vinti)", 16, BODY),
    ],
    note="I numeri si aggiornano da soli man mano che crei lead e li sposti tra gli stati.")

# ── 6. Lead elenco ───────────────────────────────────────────────────────────
content_slide(
    "I Lead: l'elenco dei contatti", "Lead",
    [
        b("La sezione “Lead” raccoglie tutti i tuoi contatti in una tabella ordinata.", 16, BODY, bullet=0),
        b("Cerca un nominativo con la barra di ricerca", 16, BODY),
        b("Filtra per Provenienza (come ti è arrivato il contatto) e per stato", 16, BODY),
        b("Vedi a colpo d'occhio nome, contatto, stato e provenienza", 16, BODY),
        b("Clicca su una riga per aprire la scheda completa del lead", 16, BODY),
    ])

# ── 7. Creare lead ───────────────────────────────────────────────────────────
content_slide(
    "Creare un nuovo lead", "Lead",
    [
        b("Premi “Nuovo lead” e compila i dati principali:", 16, BODY, bullet=0),
        b("Nome e cognome, email e telefono", 16, BODY),
        b("Provenienza del contatto (es. passaparola, sito, campagna)", 16, BODY),
        b("Capitale: puoi indicare una fascia oppure l'importo esatto", 16, BODY),
        b("Salva: il nuovo lead entra subito nella pipeline nel primo stato", 16, BODY),
    ],
    note="I campi obbligatori sono segnalati; se manca qualcosa il sistema te lo indica prima di salvare.")

# ── 8. Dettaglio lead ────────────────────────────────────────────────────────
two_col_slide(
    "La scheda di dettaglio del lead", "Lead",
    "Sintesi e dati", [
        b("In alto la “Sintesi”: stato, provenienza, capitale", 15, BODY),
        b("Anagrafica e contatti del lead", 15, BODY),
        b("Capitale come fascia o come importo preciso", 15, BODY),
        b("Cambia lo stato direttamente dalla scheda", 15, BODY),
    ],
    "Colonna “Attività”", [
        b("Al centro le note/attività in ordine di tempo", 15, BODY),
        b("Aggiungi una nota per ogni contatto o avanzamento", 15, BODY),
        b("Tieni la storia del rapporto in un unico posto", 15, BODY),
        b("Utile per ricordare il prossimo passo da fare", 15, BODY),
    ])

# ── 9. Pipeline ──────────────────────────────────────────────────────────────
content_slide(
    "La Pipeline (board kanban)", "Pipeline",
    [
        b("Mostra i lead come schede, raccolte in colonne per stato.", 16, BODY, bullet=0),
        b("Trascina una scheda da una colonna all'altra per cambiare stato", 16, BODY),
        b("Clicca una scheda per aprire il dettaglio del lead", 16, BODY),
        b("Filtra per Provenienza per concentrarti su un canale", 16, BODY),
        b("Gli stati “Vinta” e “Persa” chiudono il percorso del lead", 16, BODY),
    ],
    note="In alternativa al trascinamento col mouse, puoi cambiare stato anche dalla scheda di dettaglio.")

# ── 10. Config pipeline ──────────────────────────────────────────────────────
content_slide(
    "Personalizzare gli stati della pipeline", "Pipeline",
    [
        b("Da “Configurazione pipeline” decidi come lavora la tua board:", 16, BODY, bullet=0),
        b("Mostra o nascondi i singoli stati", 16, BODY),
        b("Riordina la sequenza degli stati", 16, BODY),
        b("Gli stati finali (Vinta / Persa) restano in fondo", 16, BODY),
    ],
    note="Imposta la pipeline una volta come preferisci: la userai così ogni giorno.")

# ── 11. Appuntamenti ─────────────────────────────────────────────────────────
content_slide(
    "Appuntamenti e agenda", "Agenda",
    [
        b("La sezione “Appuntamenti” raccoglie gli incontri collegati ai tuoi lead.", 16, BODY, bullet=0),
        b("Vedi gli appuntamenti e filtrali per stato", 16, BODY),
        b("Ogni appuntamento è collegato al lead di riferimento", 16, BODY),
        b("Ti aiuta a non perdere i prossimi contatti in agenda", 16, BODY),
    ],
    note="L'integrazione con calendari esterni (es. Google) è disattivata in questa configurazione.")

# ── 12. Export ───────────────────────────────────────────────────────────────
content_slide(
    "Esportare i tuoi dati", "Dati",
    [
        b("Puoi portare i tuoi lead fuori dal CRM quando ti serve:", 16, BODY, bullet=0),
        b("Esporta in Excel (XLSX) per analisi o archivio", 16, BODY),
        b("Esporta in JSON per usi tecnici o backup", 16, BODY),
        b("Trovi l'esportazione nel menu azioni (⋯) dell'elenco lead", 16, BODY),
    ])

# ── 13. Impostazioni ─────────────────────────────────────────────────────────
content_slide(
    "Impostazioni e aspetto", "Impostazioni",
    [
        b("Dalla sezione “Impostazioni” gestisci le preferenze del tuo spazio:", 16, BODY, bullet=0),
        b("Aspetto & brand: colori e personalizzazione visiva", 16, BODY),
        b("Integrazioni: connessioni a servizi esterni (se previste)", 16, BODY),
        b("Tema chiaro/scuro selezionabile in ogni momento dall'alto", 16, BODY),
    ])

# ── 14. Privacy/sicurezza ────────────────────────────────────────────────────
content_slide(
    "Privacy, cookie e sicurezza", "Sicurezza",
    [
        b("CustomerSpeed è pensato nel rispetto del GDPR:", 16, BODY, bullet=0),
        b("Banner cookie con scelta libera (accetta / rifiuta / preferenze)", 16, BODY),
        b("I tuoi dati sono isolati: vedi solo i lead del tuo spazio", 16, BODY),
        b("Le pagine di accesso sono protette da Google reCAPTCHA contro gli abusi", 16, BODY),
        b("Se un accesso sembra sospetto, può comparire un breve controllo “Non sono un robot”", 16, BODY),
    ],
    note="Il controllo anti-robot è normale: serve a proteggere il tuo account.")

# ── 15. Consigli + chiusura ──────────────────────────────────────────────────
s = prs.slides.add_slide(BLANK)
box(s, 0, 0, EMU_W, EMU_H, fill=INDIGO)
text(s, Inches(0.9), Inches(0.9), Inches(11.5), Inches(0.6),
     [b("In sintesi", 16, RGBColor(0xCD, 0xCE, 0xF7), bold=True, bullet=0)])
text(s, Inches(0.9), Inches(1.5), Inches(11.5), Inches(0.8),
     [b("Il tuo flusso di lavoro quotidiano", 32, WHITE, bold=True, bullet=0)])
text(s, Inches(0.95), Inches(2.7), Inches(11.4), Inches(3.2),
     [
        b("Accedi dal tuo link riservato (?org=fabio)", 18, WHITE),
        b("Controlla la Dashboard per la situazione del giorno", 18, WHITE),
        b("Inserisci i nuovi lead e tienili aggiornati con le note", 18, WHITE),
        b("Sposta i lead nella Pipeline man mano che avanzano", 18, WHITE),
        b("Esporta i dati quando ti serve un report", 18, WHITE),
     ], space_after=12)
text(s, Inches(0.95), Inches(6.4), Inches(11.4), Inches(0.6),
     [b("Buon lavoro con CustomerSpeed.", 16, RGBColor(0xE7, 0xE7, 0xFB), italic=True, bullet=0)])

prs.save("Guida-CustomerSpeed-Fabio.pptx")
print("OK:", len(prs.slides.__iter__.__self__._sldIdLst), "slide")
print("File: Guida-CustomerSpeed-Fabio.pptx")
